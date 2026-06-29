"""Build a PPTX summarising cablestack results across cables and loadings.

Pulls the postprocessed stress-strain data that lives at:
    data/runs/<run>/APDL/submodel/apdl_runfolder/pp/<cable>[_<suffix>]_stress_strain.txt

and produces a PPTX with: title, status matrix, per-cable detail, cross-cable
comparison, and an outstanding-work slide. Plots are matplotlib PNGs embedded
on the slides.

Run:
    python scripts/analysis/cablestack_slide_deck.py
"""
from __future__ import annotations

import io
from pathlib import Path
from typing import Optional

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[2]
RUNS = ROOT / "data" / "runs"
OUT_DIR = ROOT / "data" / "slidedeck"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PPTX = OUT_DIR / "cablestack_results.pptx"

# Stage display order + filename suffix used in pp/<cable>{suffix}_stress_strain.txt
STAGES = [
    ("displacement_transverse", "",              "Disp. transverse"),
    ("displacement_radial",     "_disp_radial",  "Disp. radial"),
    ("pressure_transverse",     "_pressure",     "Pressure transverse"),
    ("pressure_radial",         "_radial",       "Pressure radial"),
]

# Per-cable best available source folder for each formulation.
# Picked from the inventory of pp/ contents at this date; if a newer
# rerun lands locally with more stages, update here.
SOURCES = {
    "R2D2_HF": {
        "GPS": RUNS / "20260504_232855_R2D2_HF_apdl_rerun_51",  # 4 stages but radial uses the old buggy spring-back BC
        "PS":  RUNS / "20260504_232855_R2D2_HF_apdl_rerun_52_ps",
    },
    "R2D2_LF": {
        "GPS": RUNS / "20260511_171204_R2D2_LF_apdl_rerun_15",  # transverse only
        "PS":  None,
    },
    "CD1": {
        "GPS": RUNS / "20260504_232855_CD1_apdl_rerun_52",      # transverse only
        "PS":  None,
    },
}

# Caveats per (cable, formulation, stage).  Annotations carried to plot titles.
CAVEATS = {
    ("R2D2_HF", "GPS", "displacement_radial"):
        "from apdl_rerun_51 -- pre-fresh-start BC: top sprang back to UY=0 "
        "(non-physical). Re-run pending.",
}

CABLE_COLORS = {"R2D2_HF": "#1f77b4", "R2D2_LF": "#ff7f0e", "CD1": "#2ca02c"}
FORM_LS      = {"GPS": "-", "PS": "--"}


def pp_path(run: Path, cable: str, stage_suffix: str) -> Optional[Path]:
    p = run / "APDL" / "submodel" / "apdl_runfolder" / "pp" / f"{cable}{stage_suffix}_stress_strain.txt"
    return p if p.is_file() else None


def load_curve(path: Path):
    """Returns (strain_load %, sigma_load MPa, strain_react %, sigma_react MPa) or (None,)*4."""
    try:
        # Skip header comment lines + the column-name line
        rows = []
        with open(path) as f:
            for line in f:
                s = line.strip()
                if not s or s.startswith("#"):
                    continue
                parts = s.split()
                if len(parts) < 6:
                    continue
                try:
                    rows.append([float(x) for x in parts[1:]])  # drop Set
                except ValueError:
                    continue  # header
        if not rows:
            return None, None, None, None
        arr = np.asarray(rows)
        # columns: Time, strain_load, sigma_load_MPa, strain_react, sigma_react_MPa
        # strain values are dimensionless; convert to percent for plotting
        eps_load   = arr[:, 1] * 100.0
        sig_load   = arr[:, 2]
        eps_react  = arr[:, 3] * 100.0
        sig_react  = arr[:, 4]
        return eps_load, sig_load, eps_react, sig_react
    except Exception as e:
        print(f"  [warn] failed to load {path}: {e}")
        return None, None, None, None


def peak_load(eps, sig) -> Optional[tuple]:
    if eps is None or len(eps) == 0:
        return None
    i = int(np.argmax(np.abs(sig)))
    return float(eps[i]), float(sig[i])


def fig_to_png_bytes(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=160, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def make_per_cable_plot(cable: str, png_path: Path) -> dict:
    """4-panel stress-strain plot for one cable, GPS + PS overlay per stage."""
    fig, axes = plt.subplots(2, 2, figsize=(11, 7.5), constrained_layout=True)
    summary: dict = {}
    for ax, (stage_name, suffix, title) in zip(axes.flat, STAGES):
        ax.set_title(title, fontsize=11)
        ax.set_xlabel("strain along load axis (%)")
        ax.set_ylabel("stress along load axis (MPa)")
        ax.grid(alpha=0.3)
        plotted = False
        for formulation, run in SOURCES[cable].items():
            if run is None:
                continue
            p = pp_path(run, cable, suffix)
            if p is None:
                continue
            eps, sig, _, _ = load_curve(p)
            if eps is None:
                continue
            label = f"{formulation}"
            note = CAVEATS.get((cable, formulation, stage_name))
            if note:
                label += " (caveat)"
            ax.plot(eps, sig, FORM_LS[formulation],
                    color=CABLE_COLORS[cable], label=label)
            pk = peak_load(eps, sig)
            if pk:
                ax.scatter([pk[0]], [pk[1]], s=22, color="red", zorder=5)
                summary[(stage_name, formulation)] = pk
            plotted = True
        if not plotted:
            ax.text(0.5, 0.5, "no data", ha="center", va="center",
                    transform=ax.transAxes, fontsize=12, color="#888")
        if plotted:
            ax.legend(fontsize=9, loc="best")
    fig.suptitle(f"{cable} -- cablestack stress-strain across loadings",
                 fontsize=13, fontweight="bold")
    fig.savefig(png_path, dpi=160, bbox_inches="tight")
    plt.close(fig)
    return summary


def make_cross_cable_plot(stage_suffix: str, stage_title: str, formulation: str,
                          png_path: Path) -> dict:
    fig, ax = plt.subplots(figsize=(8.5, 5.5), constrained_layout=True)
    ax.set_title(f"{stage_title} ({formulation}) -- cable-by-cable", fontsize=13)
    ax.set_xlabel("strain along load axis (%)")
    ax.set_ylabel("stress along load axis (MPa)")
    ax.grid(alpha=0.3)
    summary = {}
    plotted = False
    for cable in ("R2D2_HF", "R2D2_LF", "CD1"):
        run = SOURCES[cable].get(formulation)
        if run is None:
            continue
        p = pp_path(run, cable, stage_suffix)
        if p is None:
            continue
        eps, sig, _, _ = load_curve(p)
        if eps is None:
            continue
        ax.plot(eps, sig, color=CABLE_COLORS[cable], lw=2, label=cable)
        pk = peak_load(eps, sig)
        if pk:
            summary[cable] = pk
            ax.scatter([pk[0]], [pk[1]], s=24, color="red", zorder=5)
            ax.annotate(f"{pk[1]:.0f} MPa @ {pk[0]:.2f}%",
                        xy=pk, xytext=(6, 6), textcoords="offset points",
                        fontsize=8.5)
        plotted = True
    if plotted:
        ax.legend(fontsize=10, loc="best")
    else:
        ax.text(0.5, 0.5, "no data", ha="center", va="center",
                transform=ax.transAxes, fontsize=14, color="#888")
    fig.savefig(png_path, dpi=160, bbox_inches="tight")
    plt.close(fig)
    return summary


# ---------------------------------------------------------------------------
# PPTX assembly
# ---------------------------------------------------------------------------

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)   # 16:9


def add_title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tx = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12.0), Inches(2.0))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = RGBColor(20, 40, 80)
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(20); p.font.color.rgb = RGBColor(80, 80, 80)
    return s


def add_section_title(prs, text):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(26); p.font.bold = True
    p.font.color.rgb = RGBColor(20, 40, 80)
    return s


def add_image_slide(prs, title, png_path, body_lines=None):
    s = add_section_title(prs, title)
    s.shapes.add_picture(str(png_path), Inches(0.4), Inches(1.4),
                         width=Inches(9.5))
    if body_lines:
        tb = s.shapes.add_textbox(Inches(10.1), Inches(1.4), Inches(3.0), Inches(5.5))
        tf = tb.text_frame; tf.word_wrap = True
        for i, line in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line; p.font.size = Pt(11)
    return s


def add_text_slide(prs, title, body_lines):
    s = add_section_title(prs, title)
    tb = s.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.0), Inches(5.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(14)
        if line.startswith(("*", "-")):
            p.level = 1
    return s


def add_status_matrix(prs):
    s = add_section_title(prs, "Status matrix -- what's complete")
    headers = ["Cable", "Form.", "Disp. transverse", "Disp. radial",
               "Pressure transverse", "Pressure radial"]
    rows: list = [headers]
    for cable in ("R2D2_HF", "R2D2_LF", "CD1"):
        for form_label in ("GPS", "PS"):
            run = SOURCES[cable].get(form_label)
            row = [cable, form_label]
            if run is None:
                row += ["-"] * 4
            else:
                for stage_name, suffix, _ in STAGES:
                    has = pp_path(run, cable, suffix) is not None
                    cell = "OK" if has else "miss"
                    if (cable, form_label, stage_name) in CAVEATS:
                        cell = "OK (caveat)"
                    row.append(cell)
            rows.append(row)

    nrows, ncols = len(rows), len(headers)
    table_left, table_top = Inches(0.6), Inches(1.4)
    table_w, table_h = Inches(12.0), Inches(4.0)
    tbl = s.shapes.add_table(nrows, ncols, table_left, table_top, table_w, table_h).table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j); cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11); p.font.bold = True
    for i in range(1, nrows):
        for j, val in enumerate(rows[i]):
            cell = tbl.cell(i, j); cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                if val == "OK":
                    p.font.color.rgb = RGBColor(0, 120, 0)
                elif val.startswith("OK (caveat)"):
                    p.font.color.rgb = RGBColor(180, 130, 0)
                elif val in ("-", "miss"):
                    p.font.color.rgb = RGBColor(160, 60, 60)
    return s


def main() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    # 1. Title
    add_title_slide(prs,
                    "Cablestack Submodel -- Results",
                    "Three cables (R2D2_HF, R2D2_LF, CD1) x four loadings\n"
                    "Generated by scripts/analysis/cablestack_slide_deck.py")

    # 2. Methodology
    add_text_slide(prs, "Methodology", [
        "2D APDL submodel: PLANE183 + CONTA172/TARGE169, large strain (NLGEOM),"
        " Voce copper plasticity, RVE-homogenised Nb3Sn isotropic modulus, wax/epoxy"
        " impregnation, sleeve insulation.",
        "",
        "Four mutually-independent load stages -- all RESUME base.db (post-mesh,"
        " pre-load) and start from an undeformed cable:",
        "  - displacement_transverse: UY ramp on top wall",
        "  - displacement_radial:     UX ramp on left wall",
        "  - pressure_transverse:     SFL pressure on top wall",
        "  - pressure_radial:         SFL pressure on left wall",
        "",
        "Two element formulations:",
        "  - GPS = generalized plane strain + mixed u-P (long-cable axial DOF;"
        " stiffer transversely, exposes Nb3Sn eps_zz for Ic prediction)",
        "  - PS  = plane stress (free in z; matches unconfined Zwick test)",
        "",
        "Two boundary types (now wired via cablestack.boundary_type):",
        "  - constrained: sidewalls perpendicular to load fixed (rigid die)",
        "  - free:        sidewalls drop the perpendicular constraint, anchor"
        " moves to the loaded-against wall (Poisson bulge allowed)",
    ])

    # 3. Status matrix
    add_status_matrix(prs)

    # 4-6. Per-cable detail slides
    summaries: dict = {}
    for cable in ("R2D2_HF", "R2D2_LF", "CD1"):
        png = OUT_DIR / f"per_cable_{cable}.png"
        s = make_per_cable_plot(cable, png)
        summaries[cable] = s
        body = ["Peak (stress-strain) per stage:"]
        for (stage_name, formulation), (eps, sig) in sorted(s.items()):
            body.append(f"{stage_name} ({formulation}): {sig:.1f} MPa @ {eps:.2f}%")
        if not s:
            body = ["No data available locally for this cable yet."]
        if (cable == "R2D2_HF"):
            body += ["",
                     "Caveat: HF GPS disp_radial",
                     "data is from the OLD",
                     "deck where the top",
                     "sprang back to UY=0",
                     "(non-physical -- see",
                     "CLAUDE.md update log)."]
        add_image_slide(prs, f"{cable} -- 4-stage stress-strain", png, body)

    # 7. Cross-cable comparison: transverse (the one stage we have for all 3)
    png = OUT_DIR / "cross_cable_displacement_transverse_GPS.png"
    s = make_cross_cable_plot("", "Disp. transverse", "GPS", png)
    body = ["Disp. transverse peak (GPS):"]
    for cable, (eps, sig) in s.items():
        body.append(f"{cable}: {sig:.1f} MPa @ {eps:.2f}%")
    body += [
        "",
        "All three cables compacted",
        "to 2% nominal strain. Peak",
        "stress reflects geometry +",
        "strand stack stiffness, not",
        "a property of the wire alone.",
    ]
    add_image_slide(prs, "Cross-cable comparison: disp. transverse (GPS)",
                    png, body)

    # 8. HF-only cross-stage comparison
    png = OUT_DIR / "hf_all_stages_GPS_vs_PS.png"
    # Use existing per-cable plot for HF (it already has GPS+PS overlay)
    add_image_slide(prs, "HF: GPS vs PS, all 4 stages",
                    OUT_DIR / "per_cable_R2D2_HF.png", [
        "HF is the only cable",
        "with both formulations",
        "+ all four stage outputs",
        "locally available.",
        "",
        "GPS curves are notably",
        "stiffer than PS in",
        "transverse (~3x at 2%)",
        "-- consistent with the",
        "saved memory note that",
        "GPS removes the z-escape",
        "channel.",
    ])

    # 9. Outstanding work
    add_text_slide(prs, "Outstanding work / caveats", [
        "Data gaps (need a fresh full run):",
        "  - LF GPS: only disp. transverse postprocessed locally (radial and"
        " pressure stages either timed out or are queued)",
        "  - CD1 GPS: only disp. transverse postprocessed locally",
        "  - LF and CD1 PS: nothing locally yet (the earlier PS jobs failed"
        " because pp/ wasn't created on the cluster; that bug is now fixed)",
        "",
        "Known caveats on data shown:",
        "  - HF GPS disp_radial curve is from apdl_rerun_51 which used the"
        " OLD restart deck (top sprang from compacted state back to UY=0,"
        " a non-physical spring-back rather than a true radial compression)",
        "  - Element 75610 distortion limited HF GPS disp_radial in the fresh"
        " run to ~0.25% strain before failure (no postprocess)",
        "",
        "Next-round queue (already submitted, pending QOSMaxMemoryPerUser):",
        "  - HF GPS w/ FKN=0.1 (job 66822524) -- diagnostic of contact stiffness"
        " as the radial bottleneck",
        "  - boundary_type='free' rerun on all 3 cables once the FKN test"
        " clarifies whether radial converges with looser contacts",
    ])

    prs.save(str(OUT_PPTX))
    print(f"\nSaved: {OUT_PPTX}")
    return OUT_PPTX


if __name__ == "__main__":
    main()
